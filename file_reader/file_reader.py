"""
KeyPhrasesExtractor extracts keyphrases from document using YAKE algorithm
Campos, R., Mangaravite, V., Pasquali, A., Jatowt, A., Jorge, A., Nunes, C. and Jatowt, A. (2020).
YAKE! Keyword Extraction from Single Documents using Multiple Local Features.
"""

import argparse
from io import BytesIO
import os

import docx
from langdetect import detect
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.pdfpage import PDFPage
import PyPDF2
from pptx import Presentation
import pandas as pd

import logging


class FilesReader(object):
    """read files from a directory"""

    def __init__(self):
        """
        Initialize the document
        """
        self.filenames = {}
        self.file_to_text = {}
        self.ext_to_funct = {".pdf": FilesReader.read_pdf,
                             ".txt": FilesReader.read_txt,
                             ".pptx": FilesReader.read_pptx,
                             ".docx": FilesReader.read_docx,
                             ".xlsx": FilesReader.read_xlsx}

        self.NO_TAG = "NO_TAG"

    @staticmethod
    def read_pdf(filename):
        """
        Read text from a pdf file.
        """
        try:
            texts = []
            with open(filename, 'rb') as outfile:
                pdfReader = PyPDF2.PdfFileReader(outfile)
                for i in range(pdfReader.numPages):
                    texts.append(pdfReader.getPage(i).extractText())
            return texts
        except:
            print("we can't read %s" % filename)
        return texts

    @staticmethod
    def read_txt(filename):
        """
        Reads a list of text
        """
        try:
            with open(filename) as f:
                texts = f.readlines()
            return texts
        except:
            print("we can't read %s" % filename)
            return []

    @staticmethod
    def read_docx(filename):
        """
        Read a list of a docx file.
        """
        try:
            doc = docx.Document(filename)
            texts = []
            for para in doc.paragraphs:
                texts.append(para.text)
            return texts
        except:
            print("we can't read %s" % filename)
            return []

    @staticmethod
    def read_pptx(filename):
        """
        Read text files from a text file.
        """
        try:
            prs = Presentation(filename)
            texts = []
            for slide in prs.slides:
                text_runs = []
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            texts.append(run.text)
            return texts
        except:
            print("we can't read %s" % filename)
            return []

    @staticmethod
    def read_xlsx(filename):
        """
        Reads an excel excel dataframe
        """
        try:
            excel_df = pd.read_excel(filename)
            texts = excel_df.to_string().split("  ")
            return [text for text in texts if text != ""]
        except:
            print("we can't read %s" % filename)
            return []

    def link_tag_to_document(self, filename, tag):
        """
        Add the document to the document.
        """
        if filename in list(self.filenames.keys()):
            self.filenames[filename]["tag"] = tag

    def link_tag_to_documents(self, filenames, tag):
        """
        Convert documents to a list of documents.
        """
        for filename in filenames:
            self.link_tag_to_document(filename, tag)

    def unlink_tag_from_document(self, filename):
        """
        Removes an html tag from a file.
        """
        if filename in list(self.filenames.keys()):
            self.filename[filename]["tag"] = self.NO_TAG

    def read_documents(self, filenames):
        """
        Read documents to the documents to - many documents.
        """
        for filename in filenames:
            ext = os.path.splitext(filename)[1]
            read_ext = self.ext_to_funct.get(ext, None)
            if read_ext is not None:
                texts = read_ext(filename)
                texts_str = " ".join(texts)
                if len(texts_str) > 20:
                    lang = detect(texts_str)
                    self.filenames.setdefault(filename, {})
                    self.filenames[filename]["lang"] = lang
                    self.filenames[filename]["tag"] = self.NO_TAG
                    self.filenames[filename]["text"] = texts
                    self.filenames[filename]["nb_caractere"] = len(texts_str)

    def remove_documents(self, filenames):
        """
        Remove documents from the documents.
        """
        for filename in filenames:
            self.filenames.pop(filename, None)

    def get_documents(self):
        """
        Return a list of documents.
        """
        documents = {}
        for filename, file_data in self.filenames.items():
            lang = file_data["lang"]
            tag = file_data["tag"]
            text = file_data["text"]

            documents.setdefault(lang, {})
            documents[lang].setdefault(tag, {})

            documents[lang][tag][filename] = text

        return documents

    def extract_infos_from_file(self, filename):
        """
        Extracts infos from a file.
        """
        ext = os.path.splitext(filename)[1]
        read_ext = self.ext_to_funct.get(ext, None)
        if read_ext is not None:
            texts = read_ext(filename)

            texts_str = " ".join(texts)
            if len(texts_str) > 20:
                lang = detect(texts_str)
                nb_words = len(texts_str.split(" "))
                return {"lang": lang, "nb_words": nb_words}
        return None
